import pandas as pd
import re
import os

distance = 5 #目標基因的上下n個
score = 2 #出現次數

class Chromosome: #創立一個chromosome 的物件
    count = 0 #總共有染色體數量
    specieslist = []     #宣告空的物種列表
    def __init__(self, filename):  #initilize filename
        self.filename = filename
        Chromosome.count+=1


        def get_data(filename):  #取得網路上下載的csv檔案
            data = pd.read_csv(r"C:\Users\junji\Myworkspace\Mylabwork\data\syntenic" + "\\" + filename, skiprows=[0]) #讀取
            data["Gene symbol"] = data["Gene symbol"].apply(lambda x: x.upper()) #將symbol統一為大寫（cacng1 -> CACNG1）
            try:
                data = data.drop('Name', axis=1).reset_index().drop("index", axis = 1).dropna() #刪除csv資料内的 'name' column
                return data.dropna()
            except: #若原本就沒有 'name' column的情況
                return data

        def get_species(filename):  #利用檔案名字來取得物種名稱 （檔案名字格式一定要是物種名含空格，接著_後面為染色體名字）
            species = re.search("[A-Za-z\s]+", filename) #正則表達取得符合字母以及底綫的格式的字串
            return species.group(0) #因為符合格式的字串被存在一個串列，因此我們只需要group(0)，第一個的物件

        def get_chromosome_location(filename):  #利用檔案名字來取得染色體的名字 (Anolis carolinensis_scaffold NW_003338810.1.CSV -> scaffold NW_003338810.1)
            temp = re.sub(".CSV", "", filename)  #去除.csv
            location = re.search("[_]([A-Za-z0-9\s_.]+)", temp) #正則表達取得符合字母以及底綫的格式的字串
            return location.group(1) #為了把re格式的[_]去掉 (_scaffold NW_003338810.1 -> scaffold NW_003338810.1)

        def get_targetgenes(data):  #傳入data，取得 cacng 或 tmem114 或 235 target 基因放進一個串列内
            targetgeneslist = []  
            for gene in data["Gene symbol"]:
                if "CACNG" in gene.upper():
                    targetgeneslist.append(gene)
                elif "TMEM114" in gene.upper():
                    targetgeneslist.append(gene)
                elif "TMEM235" in gene.upper():
                    targetgeneslist.append(gene)
            return targetgeneslist


        def get_neighbour_genes(targetgenes,data,n): #取得鄰近n個基因的函數
            neighbourindex = []  #宣告一個新的鄰近基因的引數 (利用引數來取得基因的目的是因為引述可以排列)
            neighbourgene = []
            for index, gene in enumerate(data["Gene symbol"]):  #a loop to save neighbourgene index
                if gene in targetgenes:
                    for i in range (index-n, index+n+1):
                        neighbourindex.append(i)
            neighbourindex = list(dict.fromkeys(neighbourindex)) #刪掉重叠的基因引數 (因為可能cacng1，4，5很近，他們鄰近的會重複出現)
            for i in neighbourindex: #利用鄰近基因的引數來取得鄰近基因的名字
                try: #防止索引的值超出data的範圍 (如果cacng出現在最前端或最後端)
                    neighbourgene.append(data["Gene symbol"][i])
                except: #超出的話會出錯，要除錯
                    pass
            return neighbourgene

        self.data = get_data(self.filename)
        self.species = get_species(self.filename)
        self.location = get_chromosome_location(self.filename)
        self.targetgenes = get_targetgenes(self.data)
        self.neighbour = get_neighbour_genes(self.targetgenes,self.data,distance)
        self.textx = None
        self.texty = None

        if self.species not in Chromosome.specieslist:
            Chromosome.specieslist.append(self.species)

    def set_genes(self,gene):
        self.genes = gene


filelist = os.listdir(r"C:\Users\junji\Myworkspace\Mylabwork\data\syntenic")  #get file list
for index, file in enumerate(filelist):  #a loop to initialize class chromosome
    exec(f"c{index} = Chromosome(file)")  #initialize chromosome
    exec('print(f"c{index}: ", end = "")')  #print initialized class name
    exec(f"print(c{index}.species, c{index}.location, c{index}.neighbour)") #print species and location
    print()

print(filelist)

speciesdict = {}        #############create a speciesdict
for i in Chromosome.specieslist:
    speciesdict[i] = []

for index in range(Chromosome.count):  #a loop to initialize class chromosome
    exec(f"speciesdict[c{index}.species].append('c{index}')")

df = pd.DataFrame()

for index, file in enumerate(filelist):  #add column name species and location to class.data (Dataframe)
     exec(f"df{index} = c{index}.data")  #assign df 
     exec(f"df{index}.insert(0,'Species', c{index}.species)")  #insert species
     exec(f"df{index}.insert(1,'Location', c{index}.location)")  #insert location
     exec(f"df = df.append(df{index}, ignore_index=True)")



distance = distance  #distance as neighbour
neighbourindex = []  #a new list to save index of neighbour genes
for index, gene in enumerate(df["Gene symbol"]):  #a loop to save neighbourgene index
    if ("CACNG") in gene:
        for i in range (index-distance, index+distance+1):
            neighbourindex.append(i)
    elif ("TMEM114") in gene:
        for i in range (index-distance, index+distance+1):
            neighbourindex.append(i)
    elif ("TMEM235") in gene:
        for i in range (index-distance, index+distance+1):
            neighbourindex.append(i)

neighbourindex = sorted(list(set(neighbourindex)))  #remove duplicates
newdf = df.iloc[neighbourindex, :-1]  #a new data by neighbourindex and remove name column


scoredict = {}  #a dictionary to store score
for gene in newdf["Gene symbol"]:  #loop to calculate score
    if "CACNG" in gene:
        scoredict[gene] = 100
    elif "TMEM114" in gene:
        scoredict[gene] = 100
    elif "TMEM235" in gene:
        scoredict[gene] = 100
    else:
        try:
            scoredict[gene]
            scoredict[gene].append(1)  #add score
        except:
            scoredict[gene] = [1]  #build score

for gene in scoredict:  #append list item to score
    if type(scoredict[gene]) == list:
        scoredict[gene] = len(scoredict[gene])


score = score  ###set a minimum score
genelist = []  #gene that higher than score
for gene in scoredict: 
    if scoredict[gene] >= score:
        genelist.append(gene)


for index in range (Chromosome.count):  #remove each neighbourgene that not in genelist(score lower than n)
    templist = []
    exec(f"""
for neighbourgene in c{index}.neighbour:
    if neighbourgene.upper() in genelist:
        templist.append(neighbourgene.upper())""")  #append score that is higher than n to a templist
    exec(f"c{index}.set_genes(templist)")  #set Chromosome.gene as templist
    exec(f'print("c{index}: ", end="")')  #print 
    exec(f"print(c{index}.species, c{index}.location, c{index}.genes)") #print species and location
    print()


buildxy = """
initialpoint = 300
tempx = initialpoint
tempy = initialpoint
step = 80
species = c0.species
c0.textx = initialpoint
c0.texty = initialpoint
for index in range(1,Chromosome.count):
    exec(f'''
if c{index}.species == species:
    c{index}.textx = tempx
    c{index}.texty = c{index-1}.texty + step*len(c{index-1}.genes) + 300
else:
    tempx += 1000
    species = c{index}.species
    c{index}.texty = initialpoint
    c{index}.textx = tempx

#print(c{index}.textx,c{index}.texty)''') 
"""

exec(buildxy)



draw = """
from PIL import Image, ImageDraw, ImageFont

exec(f"newImage = Image.new('RGBA', (c{Chromosome.count-1}.textx + 800,10000), 'White')")
drawObj = ImageDraw.Draw(newImage)

fontInfo1 = ImageFont.truetype('C:\Windows\Fonts\calibri.ttf', 80)
fontInfo2 = ImageFont.truetype('C:\Windows\Fonts\calibri.ttf', 50)

for i in range(Chromosome.count):
    exec(f'drawObj.text((c{i}.textx, c{i}.texty), c{i}.location, fill="black", font=fontInfo2)')
    exec(f'''
ladder = 0
step = 80
for gene in c{i}.genes:
    if c{i}.texty == initialpoint:
        drawObj.text((c{i}.textx, 50), c{i}.species, fill="black",font=fontInfo1)
    ladder += step
    if "CACNG" in gene:
        drawObj.text((c{i}.textx, c{i}.texty + ladder), gene, fill="red", font=fontInfo2)
    elif "TMEM114" in gene:
        drawObj.text((c{i}.textx, c{i}.texty + ladder), gene, fill="red", font=fontInfo2)
    elif "TMEM235" in gene:
        drawObj.text((c{i}.textx, c{i}.texty + ladder), gene, fill="red", font=fontInfo2)
    else:
        drawObj.text((c{i}.textx, c{i}.texty + ladder), gene, fill="black", font=fontInfo2)
    ''')
newImage.show()
# newImage.save("out.png")
"""

exec(draw)


sortingspecies = [
    "Petromyzon marinus",
    "Callorhinchus milii", 
    "Danio rerio", 
    "Oreochromis niloticus", 
    "Xenopus tropicalis", 
    "Chrysemys picta", 
    "Anolis carolinensis", 
    "Canis lupus familiaris",
    "Mus musculus",
    "Homo sapiens", 
    "Gallus gallus"
    ]

sortedlist = []
tups = list(speciesdict.items())
tups[1][0]
for index,i in enumerate(sortingspecies):
    for index2, j in enumerate(tups):
        if j[0] == i:
            sortedlist.append(j)

speciesdict = dict(sortedlist)


from PIL import Image, ImageDraw, ImageFont

fontInfo1 = ImageFont.truetype('C:\Windows\Fonts\calibri.ttf', 60)
fontInfo2 = ImageFont.truetype('C:\Windows\Fonts\calibri.ttf', 50)

for index in range (Chromosome.count):

    exec(f"""newImage{index} = Image.new('RGBA', (850, 80*len(c{index}.genes)+250), 'white')
    
drawObj = ImageDraw.Draw(newImage{index})

ladder = 10

drawObj.text((10, ladder), c{index}.location, fill="black", font=fontInfo1)

ladder += 100

for gene in c{index}.genes:
    if "CACNG" in gene:
        drawObj.text((10, ladder), gene, fill="red", font=fontInfo2)
    elif "TMEM114" in gene:
        drawObj.text((10, ladder), gene, fill="red", font=fontInfo2)
    elif "TMEM235" in gene:
        drawObj.text((10, ladder), gene, fill="red", font=fontInfo2)
    elif scoredict[gene] >= 7:
        drawObj.text((10, ladder), gene, fill="blue", font=fontInfo2)
    else:

        drawObj.text((10, ladder), gene, fill="black", font=fontInfo2)
    ladder += 80

c{index}.image = newImage{index}

newImage{index}.save(r'C:\\Users\\junji\\Myworkspace\\Mylabwork\\data\\image\\C{index}.png')

""")


import pptx
from pptx.util import Inches, Pt

prs = pptx.Presentation()
prs.slide_width = Inches(11)
prs.slide_height = Inches(6)

slide_1 = prs.slides.add_slide(prs.slide_layouts[6])

lad1 = 0
for species in speciesdict:
    tb = slide_1.shapes.add_textbox(Inches(lad1-0.08),Inches(0),Inches(1),Inches(0.25))    #textbox
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = species
    p.font.size = Pt(7)
    lad2 = 0.25
    for c in speciesdict[species]:
        exec(f"tb = slide_1.shapes.add_picture(r'C:\\Users\\junji\\Myworkspace\\Mylabwork\\data\\image\\{c}.png',Inches(lad1),Inches(lad2),width=Inches(1))")
        exec(f"lad2 += (80*len({c}.genes)+250)/850")
    lad1+=1

prs.save(r"C:\Users\junji\Desktop\test.pptx")